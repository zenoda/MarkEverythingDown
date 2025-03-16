from processors.base import BaseDocumentProcessor, StructuredDocument, DocumentSection, DocumentElement, DocumentType
from pathlib import Path
from pptx import Presentation

class PowerPointProcessor(BaseDocumentProcessor):
    """Processor for PowerPoint (PPTX) documents"""
    
    def process(self, file_path: str) -> StructuredDocument:
        """Process PowerPoint presentation"""
        document = StructuredDocument(
            title=Path(file_path).stem,
            source_file=file_path,
            doc_type=DocumentType.POWERPOINT
        )
        
        try:
            # Open presentation
            presentation = Presentation(file_path)
            
            # Storage for markdown output
            markdown_parts = []
            markdown_parts.append(f"# {Path(file_path).stem}\n\n")
            
            # Process each slide
            for i, slide in enumerate(presentation.slides):
                # Create section for each slide
                slide_section = DocumentSection(title=f"Slide {i+1}", level=1)
                slide_section.metadata["slide_number"] = i+1
                
                # Add slide title to markdown
                markdown_parts.append(f"## Slide {i+1}\n\n")
                
                # Process slide title if available
                if slide.shapes.title:
                    title_text = slide.shapes.title.text
                    markdown_parts.append(f"### {title_text}\n\n")
                    
                    title_element = DocumentElement(
                        content=title_text,
                        element_type="heading",
                        metadata={"level": 3}
                    )
                    slide_section.add_element(title_element)
                
                # Process text elements
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip() and shape != slide.shapes.title:
                        shape_text = shape.text.strip()
                        # Skip if it's same as title
                        if slide.shapes.title and shape_text == slide.shapes.title.text:
                            continue
                        
                        markdown_parts.append(f"{shape_text}\n\n")
                        
                        # Add text as paragraph
                        text_element = DocumentElement(
                            content=shape_text,
                            element_type="paragraph"
                        )
                        slide_section.add_element(text_element)
                
                # Add notes if available
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    markdown_parts.append(f"**Notes:**\n\n{notes_text}\n\n")
                    
                    notes_element = DocumentElement(
                        content=notes_text,
                        element_type="paragraph",
                        metadata={"is_notes": True}
                    )
                    slide_section.add_element(notes_element)
                
                # Add divider between slides
                markdown_parts.append("---\n\n")
                
                # Add section to document
                document.add_section(slide_section)
            
            # Store combined markdown
            document.metadata["markdown"] = "".join(markdown_parts)
            
        except Exception as e:
            print(f"Error processing PowerPoint: {e}")
            error_section = DocumentSection(title="Error")
            error_section.add_element(DocumentElement(
                content=f"Failed to process PowerPoint: {str(e)}",
                element_type="paragraph"
            ))
            document.add_section(error_section)
        
        return document